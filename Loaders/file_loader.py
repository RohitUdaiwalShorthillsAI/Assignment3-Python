import os
import docx
import PyPDF2
from pptx import Presentation

class FileLoader:
    @staticmethod
    def load_file(file_path):
        file_extension = os.path.splitext(file_path)[1].lower()

        try:
            # Check if the file is a PDF
            if file_extension == '.pdf':
                # Load and extract content from PDF
                with open(file_path, 'rb') as file:
                    reader = PyPDF2.PdfReader(file)
                    content = ""

                    # Check if the PDF is encrypted
                    if reader.is_encrypted:
                        if reader.decrypt("") == 0:
                            return "Error loading PDF: The file is password-protected."

                    # Extract text from each page
                    for page in range(len(reader.pages)):
                        content += reader.pages[page].extract_text()

                return content

            # Check if the file is a DOCX
            elif file_extension == '.docx':
                # Load and extract content from DOCX
                doc = docx.Document(file_path)
                content = "\n".join([para.text for para in doc.paragraphs])
                return content

            # Check if the file is a PPTX
            elif file_extension == '.pptx':
                # Load and extract content from PPTX
                presentation = Presentation(file_path)
                content = ""

                # Extract text from each slide
                for slide in presentation.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            content += shape.text + "\n"

                return content

            else:
                return "Unsupported file format."

        except Exception as e:
            return f"Error loading file: {e}"
