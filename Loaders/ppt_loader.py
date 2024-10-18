from pptx import Presentation  # To handle PPT files
from .abstract_loader import FileLoader  # Import the abstract FileLoader class
import os  # To handle file path operations

class PPTLoader(FileLoader):
    """
    A concrete class for loading and processing PPTX files.

    This class inherits from the abstract FileLoader class and provides
    implementation for validating and loading PPTX files.

    Methods:
        validate_file(file_path): 
            Validates if the file exists and is a valid PPTX file.

        load_file(file_path): 
            Loads the PPTX file and extracts its text content from slides.
    """

    def validate_file(self, file_path):
        """
        Validate if the provided file exists and if it is a valid PPTX file.

        Parameters:
            file_path (str): The path to the PPTX file to be validated.

        Returns:
            tuple: A tuple containing:
                - A boolean value indicating if the file is valid (True) or not (False).
                - A string message providing additional information about the validation result.
        
        Example:
            is_valid, message = ppt_loader.validate_file('path/to/file.pptx')
            if is_valid:
                print("File is valid")
            else:
                print(message)
        """
        if not os.path.exists(file_path):
            return False, "File does not exist."  # File doesn't exist
        if not file_path.endswith('.pptx'):
            return False, "Invalid file type. Expected a PPTX file."  # Wrong file type
        return True, "File is valid."  # File exists and is valid

    def load_file(self, file_path):
        """
        Load the PPTX file and extract its text content.

        Iterates through each slide in the presentation and extracts text
        from all shapes that have a text attribute.

        Parameters:
            file_path (str): The path to the PPTX file to be loaded and processed.

        Returns:
            str: The extracted text content from the PPTX file or an error message if any issue occurs.
        
        Raises:
            Exception: If there is an error during file loading.

        Example:
            content = ppt_loader.load_file('path/to/file.pptx')
            print(content)
        """
        try:
            # Open the PPTX file
            presentation = Presentation(file_path)
            content = ""

            # Iterate through each slide in the presentation
            for slide in presentation.slides:
                # Iterate through each shape in the slide
                for shape in slide.shapes:
                    # Check if the shape has text and add it to the content
                    if hasattr(shape, "text"):
                        content += shape.text + "\n"  # Extract text from the shape

            return content  # Return the extracted content
        except Exception as e:
            # Return an error message if something goes wrong
            return f"Error loading PPT: {e}"
