import os
import io
import fitz  # PyMuPDF for PDF handling
import docx  # python-docx for DOCX handling
import pptx  # python-pptx for PPTX handling
import sqlite3
import pdfplumber
from PyPDF2 import PdfReader
from PIL import Image
from io import BytesIO
from docx import Document
from pptx import Presentation

class FileDataExtractor:
    """
    A unified class for extracting text, images, tables, and links from PDF, DOCX, and PPTX files.
    
    Methods:
        extract_text(): Extracts text content from the file.
        extract_images(): Extracts image data from the file.
        extract_tables(): Extracts tables from the file.
        extract_links(): Extracts hyperlinks from the file.
    """

    def __init__(self, file_path):
        self.file_path = file_path
        self.file_extension = os.path.splitext(file_path)[1].lower()

    def extract_text(self):
        """
        Extract text content from the file based on its type (PDF, DOCX, PPTX).
        """
        if self.file_extension == '.pdf':
            with open(self.file_path, "rb") as f:
                reader = PdfReader(f)
                text = "".join([page.extract_text() or "" for page in reader.pages])
                metadata = self._extract_metadata(reader.metadata)
            return text, metadata

        elif self.file_extension == '.docx':
            doc = Document(self.file_path)
            text = "\n".join([para.text for para in doc.paragraphs])
            metadata = self._extract_metadata(doc.core_properties)
            return text, metadata

        elif self.file_extension == '.pptx':
            presentation = Presentation(self.file_path)
            text = "\n".join([shape.text for slide in presentation.slides for shape in slide.shapes if hasattr(shape, "text")])
            metadata = self._extract_metadata(presentation.core_properties)
            return text, metadata

        else:
            raise ValueError("Unsupported file format. Only PDF, DOCX, and PPTX are supported.")

    def extract_images(self):
        """
        Extract images from the file based on its type (PDF, DOCX, PPTX).
        """
        images = []
        if self.file_extension == '.pdf':
            with fitz.open(self.file_path) as doc:
                for page_num in range(len(doc)):
                    for img in doc[page_num].get_images(full=True):
                        images.append(self._process_image("pdf", img, doc, page_num + 1))

        elif self.file_extension == '.docx':
            doc = Document(self.file_path)
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    image_blob = rel.target_part.blob
                    images.append(self._process_image("docx", image_blob))

        elif self.file_extension == '.pptx':
            presentation = Presentation(self.file_path)
            for slide_num, slide in enumerate(presentation.slides):
                for shape in slide.shapes:
                    if hasattr(shape, "image") and shape.image:
                        images.append(self._process_image("pptx", shape, slide_num + 1))

        else:
            raise ValueError("Unsupported file format. Only PDF, DOCX, and PPTX are supported.")

        return images

    def extract_tables(self):
        """
        Extract tables from the file based on its type (PDF, DOCX, PPTX).
        """
        tables = []
        if self.file_extension == '.pdf':
            with pdfplumber.open(self.file_path) as pdf:
                for page in pdf.pages:
                    tables.extend(page.extract_tables() or [])

        elif self.file_extension == '.docx':
            doc = Document(self.file_path)
            tables = [[self._extract_table_row(row) for row in table.rows] for table in doc.tables]

        elif self.file_extension == '.pptx':
            presentation = Presentation(self.file_path)
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if shape.has_table:
                        tables.append([self._extract_table_row(row) for row in shape.table.rows])

        else:
            raise ValueError("Unsupported file format. Only PDF, DOCX, and PPTX are supported.")

        return tables

    def extract_links(self):
        """
        Extract links from the file based on its type (PDF, DOCX, PPTX).
        """
        links = []
        if self.file_extension == '.pdf':
            with fitz.open(self.file_path) as doc:
                for page_num in range(len(doc)):
                    links.extend(self._extract_pdf_link(doc[page_num], page_num + 1))

        elif self.file_extension == '.docx':
            doc = Document(self.file_path)
            for para in doc.paragraphs:
                for rel in para._p.xpath('.//w:hyperlink'):
                    rId = rel.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
                    if rId in para.part.rels:
                        links.append({"url": para.part.rels[rId].target_ref, "page_number": None})

        elif self.file_extension == '.pptx':
            presentation = Presentation(self.file_path)
            links = [{"url": shape.hyperlink.address, "slide_number": slide_num + 1}
                     for slide_num, slide in enumerate(presentation.slides)
                     for shape in slide.shapes if hasattr(shape, "hyperlink") and shape.hyperlink.address]

        else:
            raise ValueError("Unsupported file format. Only PDF, DOCX, and PPTX are supported.")

        return links

    def _extract_metadata(self, properties):
        """
        Extract metadata from the file properties for DOCX and PPTX files.
        """
        return {
            "author": getattr(properties, 'author', ''),
            "created": getattr(properties, 'created', ''),
            "last_modified_by": getattr(properties, 'last_modified_by', ''),
            "title": getattr(properties, 'title', '')
        }


    def _process_image(self, file_type, img, doc=None, page_number=None):
        """
        Process images from PDF, DOCX, and PPTX files.
        """
        if file_type == "pdf":
            base_image = doc.extract_image(img[0])
            return {
                "image": sqlite3.Binary(base_image["image"]),
                "image_format": base_image["ext"],
                "image_resolution": f"{base_image['width']}x{base_image['height']}",
                "page_number": page_number
            }
        elif file_type == "docx":
            image_blob = img  # Directly use the blob for DOCX
            image = Image.open(BytesIO(image_blob))
            return {
                "image": sqlite3.Binary(image_blob),
                "image_format": image.format.lower(),
                "image_resolution": f"{image.width}x{image.height}",
            }
        elif file_type == "pptx":
            image_blob = io.BytesIO(img.image.blob)
            img = Image.open(image_blob)
            return {
                "image": sqlite3.Binary(image_blob.getvalue()),
                "image_format": img.format.lower(),
                "image_resolution": f"{img.width}x{img.height}",
                "slide_number": page_number
            }
        else:
            raise ValueError("Unsupported file type for image processing.")

    # --------------------------- PDF Helper Methods --------------------------- #

    def _extract_pdf_link(self, page, page_number):
        return [{"url": link.get("uri", ""), "page_number": page_number} for link in page.get_links()]

    # --------------------------- Common Helper Methods --------------------------- #

    def _extract_table_row(self, row):
        return [cell.text.strip() for cell in row.cells]
