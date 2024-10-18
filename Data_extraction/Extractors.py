import fitz  # PyMuPDF for PDF handling
import docx  # python-docx for DOCX
import pptx  # python-pptx for PPTX
import io
from PyPDF2 import PdfReader
from PIL import Image
import sqlite3
import pdfplumber
from docx import Document
from io import BytesIO

class DataExtractor:
    def __init__(self, loader):
        self.loader = loader
        self.file_path = loader.file_path


class PDFDataExtractor(DataExtractor):
    def extract_text(self):
        with open(self.file_path, "rb") as f:
            reader = PdfReader(f)
            text = "".join([page.extract_text() or "" for page in reader.pages])
            metadata = self._extract_metadata(reader.metadata)
        return text, metadata

    def extract_images(self):
        images = []
        with fitz.open(self.file_path) as doc:
            for page_num, page in enumerate(doc):
                for img in page.get_images(full=True):
                    images.append(self._process_image(img, doc, page_num + 1))
        return images

    def extract_tables(self):
        tables = []
        with pdfplumber.open(self.file_path) as pdf:
            for page in pdf.pages:
                tables.extend(page.extract_tables() or [])
        return tables

    def extract_links(self):
        links = []
        with fitz.open(self.file_path) as doc:
            for page_num, page in enumerate(doc):
                links.extend(self._extract_links(page, page_num + 1))
        return links

    def _extract_metadata(self, metadata):
        return {k.lower(): v for k, v in metadata.items() if v}

    def _process_image(self, img, doc, page_number):
        base_image = doc.extract_image(img[0])
        return {
            "image": sqlite3.Binary(base_image["image"]),
            "image_format": base_image["ext"],
            "image_resolution": f"{base_image['width']}x{base_image['height']}",
            "page_number": page_number
        }

    def _extract_links(self, page, page_number):
        return [{"url": link.get("uri", ""), "page_number": page_number} for link in page.get_links()]


class DOCXDataExtractor(DataExtractor):
    def extract_text(self):
        doc = docx.Document(self.file_path)
        text = "\n".join([para.text for para in doc.paragraphs])
        metadata = self._extract_metadata(doc.core_properties)
        return text, metadata

    def extract_images(self):
        doc = docx.Document(self.file_path)
        return [self._process_image(rel) for rel in doc.part.rels.values() if "image" in rel.target_ref]

    def extract_tables(self):
        doc = docx.Document(self.file_path)
        return [[self._extract_table_row(row) for row in table.rows] for table in doc.tables]

    def extract_links(self):
        doc = docx.Document(self.file_path)
        return [self._extract_link(para) for para in doc.paragraphs if self._extract_link(para)]

    def _extract_metadata(self, properties):
        return {
            "author": properties.author,
            "created": properties.created,
            "last_modified_by": properties.last_modified_by,
            "title": properties.title
        }

    def extract_images(self):
        # Load the DOCX file
        doc = Document(self.file_path)
        images = []

        # Iterate through all the parts (image files in the document)
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                image_blob = rel.target_part.blob  # Extract the image binary data
                
                # Open the image using PIL to get its resolution
                image = Image.open(BytesIO(image_blob))
                width, height = image.size  # Get image dimensions

                images.append({
                    "image": sqlite3.Binary(image_blob),  # Store image binary for DB
                    "image_format": rel.target_ref.split(".")[-1],  # Extract file extension
                    "image_resolution": f"{width}x{height}",  # Store resolution as WxH
                })

        return images

    def _extract_table_row(self, row):
        return [cell.text.strip() for cell in row.cells]

    def _extract_link(self, para):
        for rel in para._p.xpath('.//w:hyperlink'):
            rId = rel.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
            if rId in para.part.rels:
                return {"url": para.part.rels[rId].target_ref, "page_number": None}


class PPTDataExtractor(DataExtractor):
    def extract_text(self):
        presentation = pptx.Presentation(self.file_path)
        text = "\n".join([shape.text for slide in presentation.slides for shape in slide.shapes if hasattr(shape, "text")])
        metadata = self._extract_metadata(presentation.core_properties)
        return text, metadata

    def extract_images(self):
        presentation = pptx.Presentation(self.file_path)
        return [self._process_image(shape, slide_num + 1) for slide_num, slide in enumerate(presentation.slides)
                for shape in slide.shapes if hasattr(shape, "image") and shape.image]

    def extract_tables(self):
        presentation = pptx.Presentation(self.file_path)
        return [[self._extract_table_row(row) for row in shape.table.rows] for slide in presentation.slides for shape in slide.shapes if shape.has_table]

    def extract_links(self):
        presentation = pptx.Presentation(self.file_path)
        return [{"url": shape.hyperlink.address, "slide_number": slide_num + 1} for slide_num, slide in enumerate(presentation.slides) for shape in slide.shapes if hasattr(shape, "hyperlink") and shape.hyperlink.address]

    def _extract_metadata(self, properties):
        return {
            "author": properties.author,
            "created": properties.created,
            "last_modified_by": properties.last_modified_by,
            "title": properties.title
        }

    def _process_image(self, shape, slide_number):
        image_blob = io.BytesIO(shape.image.blob)
        img = Image.open(image_blob)
        return {
            "image": sqlite3.Binary(image_blob.getvalue()),
            "image_format": img.format.lower(),
            "image_resolution": f"{img.width}x{img.height}",
            "slide_number": slide_number
        }

    def _extract_table_row(self, row):
        return [cell.text.strip() for cell in row.cells]
