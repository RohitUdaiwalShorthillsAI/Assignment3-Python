from Loaders import PDFLoader 
from Loaders import PPTLoader
from Loaders import DOCXLoader 
from abc import ABC, abstractmethod

class DataExtractor(ABC):
    def __init__(self, file_loader):
        self.file_loader = file_loader

    @abstractmethod
    def extract_text(self):
        pass

    @abstractmethod
    def extract_links(self):
        pass

    @abstractmethod
    def extract_images(self):
        pass

    @abstractmethod
    def extract_tables(self):
        pass

from PyPDF2 import PdfReader
import re

class PDFDataExtractor(DataExtractor):
    
    def extract_text(self):
        try:
            content = ""
            metadata = []
            reader = PdfReader(self.file_loader.file_path)
            for page_num, page in enumerate(reader.pages):
                page_text = page.extract_text()
                metadata.append({
                    'page_number': page_num + 1,
                    'text': page_text
                })
                content += page_text
            return content, metadata
        except Exception as e:
            return f"Error extracting text from PDF: {e}"

    def extract_links(self):
        try:
            reader = PdfReader(self.file_loader.file_path)
            links = []
            for page_num, page in enumerate(reader.pages):
                # Simple regex to find URLs in text (pdfminer should be used for extracting true links)
                urls = re.findall(r'(https?://\S+)', page.extract_text())
                for url in urls:
                    links.append({
                        'url': url,
                        'page_number': page_num + 1
                    })
            return links
        except Exception as e:
            return f"Error extracting links from PDF: {e}"

    def extract_images(self):
        # PDF image extraction logic (can use PyMuPDF or pdfminer for detailed extraction)
        return "Image extraction not implemented yet for PDF."

    def extract_tables(self):
        # PDF table extraction logic (can use tabula-py or Camelot)
        return "Table extraction not implemented yet for PDF."

import docx

class DOCXDataExtractor(DataExtractor):
    
    def extract_text(self):
        try:
            doc = docx.Document(self.file_loader.file_path)
            content = ""
            metadata = []
            
            for i, para in enumerate(doc.paragraphs):
                para_text = para.text
                para_metadata = {
                    'paragraph_number': i + 1,
                    'text': para_text,
                    'font': para.runs[0].font.name if para.runs else None,
                    'font_size': para.runs[0].font.size.pt if para.runs and para.runs[0].font.size else None,
                    'bold': para.runs[0].font.bold if para.runs else None
                }
                metadata.append(para_metadata)
                content += para_text + "\n"
                
            return content, metadata
        except Exception as e:
            return f"Error extracting text from DOCX: {e}"

    def extract_links(self):
        try:
            doc = docx.Document(self.file_loader.file_path)
            links = []
            for para in doc.paragraphs:
                for run in para.runs:
                    if run.hyperlink:
                        links.append({
                            'text': run.text,
                            'url': run.hyperlink.target,
                            'paragraph_number': para.paragraph_number
                        })
            return links
        except Exception as e:
            return f"Error extracting links from DOCX: {e}"

    def extract_images(self):
        # DOCX image extraction logic (not implemented here, requires custom work with embedded images)
        return "Image extraction not implemented yet for DOCX."

    def extract_tables(self):
        try:
            doc = docx.Document(self.file_loader.file_path)
            tables_metadata = []
            for table_idx, table in enumerate(doc.tables):
                rows = len(table.rows)
                cols = len(table.columns)
                tables_metadata.append({
                    'table_index': table_idx + 1,
                    'rows': rows,
                    'columns': cols,
                    'table_dimensions': f"{rows}x{cols}"
                })
            return tables_metadata
        except Exception as e:
            return f"Error extracting tables from DOCX: {e}"

from pptx import Presentation

class PPTDataExtractor(DataExtractor):

    def extract_text(self):
        try:
            prs = Presentation(self.file_loader.file_path)
            content = ""
            metadata = []
            for slide_num, slide in enumerate(prs.slides):
                slide_text = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text += shape.text + "\n"
                        metadata.append({
                            'slide_number': slide_num + 1,
                            'text': shape.text
                        })
                content += slide_text
            return content, metadata
        except Exception as e:
            return f"Error extracting text from PPT: {e}"

    def extract_links(self):
        try:
            prs = Presentation(self.file_loader.file_path)
            links = []
            for slide_num, slide in enumerate(prs.slides):
                for shape in slide.shapes:
                    if hasattr(shape, "text") and hasattr(shape, "hyperlink"):
                        links.append({
                            'slide_number': slide_num + 1,
                            'linked_text': shape.text,
                            'url': shape.hyperlink.address
                        })
            return links
        except Exception as e:
            return f"Error extracting links from PPT: {e}"

    def extract_images(self):
        try:
            prs = Presentation(self.file_loader.file_path)
            images_metadata = []
            for slide_num, slide in enumerate(prs.slides):
                for shape in slide.shapes:
                    if 'picture' in shape.shape_type._member_names_:
                        image = shape.image
                        images_metadata.append({
                            'slide_number': slide_num + 1,
                            'image_resolution': f"{image.width}x{image.height}",
                            'image_format': image.ext
                        })
            return images_metadata
        except Exception as e:
            return f"Error extracting images from PPT: {e}"

    def extract_tables(self):
        try:
            prs = Presentation(self.file_loader.file_path)
            tables_metadata = []
            for slide_num, slide in enumerate(prs.slides):
                for shape in slide.shapes:
                    if shape.has_table:
                        table = shape.table
                        rows = len(table.rows)
                        cols = len(table.columns)
                        tables_metadata.append({
                            'slide_number': slide_num + 1,
                            'rows': rows,
                            'columns': cols,
                            'table_dimensions': f"{rows}x{cols}"
                        })
            return tables_metadata
        except Exception as e:
            return f"Error extracting tables from PPT: {e}"

# Example usage
pdf_loader = PDFLoader()
pdf_loader.file_path = 'Document.pdf'  # PDF file path
pdf_extractor = PDFDataExtractor(pdf_loader)

# Extract text and metadata from PDF
pdf_text, pdf_text_metadata = pdf_extractor.extract_text()
print("PDF Text:", pdf_text)
print("PDF Metadata:", pdf_text_metadata)

# Similarly for DOCX and PPT
docx_loader = DOCXLoader()
docx_loader.file_path = 'Document.docx'  # DOCX file path
docx_extractor = DOCXDataExtractor(docx_loader)

# Extract tables from DOCX
docx_tables_metadata = docx_extractor.extract_tables()
print("DOCX Tables Metadata:", docx_tables_metadata)

ppt_loader = PPTLoader()
ppt_loader.file_path = 'Document.pptx'  # PPTX file path
ppt_extractor = PPTDataExtractor(ppt_loader)

# Extract images from PPT
ppt_images_metadata = ppt_extractor.extract_images()
print("PPT Images Metadata:", ppt_images_metadata)
