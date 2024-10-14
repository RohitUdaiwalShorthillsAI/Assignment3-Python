from abc import ABC, abstractmethod
import os

class FileLoader(ABC):
    @abstractmethod
    def validate_file(self, file_path):
        """Validate if the file exists and is of the correct type."""
        pass

    @abstractmethod
    def load_file(self, file_path):
        """Load and process the file."""
        pass
# Import necessary libraries
import PyPDF2  # To handle PDF files
import docx    # To handle DOCX files
from pptx import Presentation  # To handle PPT files

# PDF Loader
class PDFLoader(FileLoader):
    def validate_file(self, file_path):
        if not os.path.exists(file_path):
            return False, "File does not exist."
        if not file_path.endswith('.pdf'):
            return False, "Invalid file type. Expected a PDF file."
        return True, "File is valid."

    def load_file(self, file_path):
        try:
            with open(file_path, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                content = ""
                for page in range(len(reader.pages)):
                    content += reader.pages[page].extract_text()
            return content
        except Exception as e:
            return f"Error loading PDF: {e}"


# DOCX Loader
class DOCXLoader(FileLoader):
    def validate_file(self, file_path):
        if not os.path.exists(file_path):
            return False, "File does not exist."
        if not file_path.endswith('.docx'):
            return False, "Invalid file type. Expected a DOCX file."
        return True, "File is valid."

    def load_file(self, file_path):
        try:
            doc = docx.Document(file_path)
            content = "\n".join([para.text for para in doc.paragraphs])
            return content
        except Exception as e:
            return f"Error loading DOCX: {e}"


# PPT Loader
class PPTLoader(FileLoader):
    def validate_file(self, file_path):
        if not os.path.exists(file_path):
            return False, "File does not exist."
        if not file_path.endswith('.pptx'):
            return False, "Invalid file type. Expected a PPTX file."
        return True, "File is valid."

    def load_file(self, file_path):
        try:
            presentation = Presentation(file_path)
            content = ""
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        content += shape.text + "\n"
            return content
        except Exception as e:
            return f"Error loading PPT: {e}"

# Example usage
# pdf_loader = PDFLoader()
# docx_loader = DOCXLoader()
# ppt_loader = PPTLoader()

# # PDF File Example
# is_valid, message = pdf_loader.validate_file('Document.pdf')
# if is_valid:
#     pdf_content = pdf_loader.load_file('Document.pdf')
#     print("PDF Content:", pdf_content)
# else:
#     print(message)

# # DOCX File Example
# is_valid, message = docx_loader.validate_file('Document.docx')
# if is_valid:
#     docx_content = docx_loader.load_file('Document.docx')
#     print("DOCX Content:", docx_content)
# else:
#     print(message)

# # PPT File Example
# is_valid, message = ppt_loader.validate_file('Document.pptx')
# if is_valid:
#     ppt_content = ppt_loader.load_file('Document.pptx')
#     print("PPT Content:", ppt_content)
# else:
#     print(message)